"""Extract text from uploads using local parsers or LlamaParse."""

from __future__ import annotations

import importlib.util
import re
import unicodedata
from collections import Counter
from dataclasses import dataclass
from itertools import zip_longest
from pathlib import Path
from typing import Any

import tiktoken
from bs4 import BeautifulSoup
from llama_index.core import Document
from llama_index.core.ingestion import IngestionPipeline
from llama_index.core.node_parser import SentenceWindowNodeParser
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from .settings import get_settings
from .telemetry import wrap_outbound_call

settings = get_settings()
TOKEN_ENCODER = tiktoken.get_encoding("cl100k_base")
PDF_ARTIFACT_TYPE = "pdf_sentence_window"
PDF_WINDOW_METADATA_KEY = "window_text"
PDF_ORIGINAL_TEXT_METADATA_KEY = "original_text"
PDF_LOCAL_TOTAL_CHAR_FLOOR = 160
PDF_LOCAL_TOTAL_TOKEN_FLOOR = 40
PDF_LOCAL_SUBSTANTIVE_CHAR_FLOOR = 120
PDF_LOCAL_SUBSTANTIVE_TOKEN_FLOOR = 24
PDF_LOCAL_SUBSTANTIVE_WORD_FLOOR = 20
PDF_LOCAL_LOW_SIGNAL_CHAR_FLOOR = 60
PDF_LOCAL_LOW_SIGNAL_TOKEN_FLOOR = 18
PDF_LOCAL_LOW_SIGNAL_WORD_FLOOR = 12
PDF_LOCAL_LOW_SIGNAL_ALNUM_RATIO = 0.4
PDF_LOCAL_GARBAGE_CHAR_FLOOR = 30
PDF_LOCAL_GARBAGE_TOKEN_FLOOR = 8
PDF_LOCAL_GARBAGE_WORD_FLOOR = 6
PDF_LOCAL_GARBAGE_ALNUM_RATIO = 0.25


@dataclass(slots=True)
class PdfPage:
    page_number: int
    text: str
    removed_header_lines: int = 0
    removed_footer_lines: int = 0


@dataclass(slots=True)
class PdfExtractionResult:
    documents: list[Document]
    parse_lane: str
    page_count: int
    total_text_chars: int
    parse_diagnostics: dict[str, Any]


def _merge_spreadsheet_cell_display(data_val: Any, formula_val: Any) -> str:
    """Prefer cached results (data_only); fall back to formula / literal text.

    Excel often omits cached values for formula cells until the workbook is opened
    and recalculated. A values-only read then yields empty cells and ingestion
    reports \"0 data rows\" even though the sheet has logic in the Filters/P&L area.
    """
    if data_val is not None:
        rendered = str(data_val)
        if rendered.strip():
            return rendered
    if formula_val is not None:
        return str(formula_val).strip()
    return ""


def extract_spreadsheet_structure(path: Path) -> dict[str, Any] | None:
    suffix = path.suffix.lower()
    if suffix not in {".xlsx", ".xlsm"}:
        return None

    wb_data = load_workbook(str(path), read_only=True, data_only=True)
    wb_formula = load_workbook(str(path), read_only=True, data_only=False)
    sheets: list[dict[str, Any]] = []
    try:
        if len(wb_data.worksheets) != len(wb_formula.worksheets):
            for sheet in wb_data.worksheets:
                rows: list[list[str]] = []
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in cells):
                        rows.append(cells)
                sheets.append({"title": sheet.title, "rows": rows})
        else:
            for sheet_data, sheet_formula in zip(wb_data.worksheets, wb_formula.worksheets, strict=True):
                title = sheet_data.title
                rows: list[list[str]] = []
                for row_data, row_formula in zip_longest(
                    sheet_data.iter_rows(values_only=True),
                    sheet_formula.iter_rows(values_only=True),
                    fillvalue=None,
                ):
                    left = row_data if row_data is not None else ()
                    right = row_formula if row_formula is not None else ()
                    cells = [
                        _merge_spreadsheet_cell_display(cd, cf)
                        for cd, cf in zip_longest(left, right, fillvalue=None)
                    ]
                    if any(cell.strip() for cell in cells):
                        rows.append(cells)
                sheets.append({"title": title, "rows": rows})
    finally:
        wb_formula.close()
        wb_data.close()

    return {
        "kind": "workbook",
        "sheet_count": len(sheets),
        "sheets": sheets,
    }


def normalize_text_content(text: str) -> str:
    text = unicodedata.normalize("NFKC", text.replace("\x00", ""))
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def estimate_token_count(text: str) -> int:
    if not text.strip():
        return 0
    return len(TOKEN_ENCODER.encode(text))


def detect_section_title(text: str) -> str | None:
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("#"):
            return line.lstrip("# ").strip()[:200]
        if len(line) <= 100 and line.upper() == line and any(char.isalpha() for char in line):
            return line.title()[:200]
    return None


def _first_nonempty_line(text: str) -> str | None:
    for line in text.splitlines():
        stripped = line.strip()
        if stripped:
            return stripped
    return None


def _last_nonempty_line(text: str) -> str | None:
    for line in reversed(text.splitlines()):
        stripped = line.strip()
        if stripped:
            return stripped
    return None


def _strip_repeated_pdf_artifacts(pages: list[PdfPage]) -> list[PdfPage]:
    if len(pages) < 3:
        return pages

    threshold = max(3, int(len(pages) * 0.6))
    headers = Counter(
        line.casefold()
        for page in pages
        if (line := _first_nonempty_line(page.text)) and len(line) <= 160 and any(char.isalpha() for char in line)
    )
    footers = Counter(
        line.casefold()
        for page in pages
        if (line := _last_nonempty_line(page.text)) and len(line) <= 160 and any(char.isalpha() for char in line)
    )
    repeated_headers = {line for line, count in headers.items() if count >= threshold}
    repeated_footers = {line for line, count in footers.items() if count >= threshold}
    if not repeated_headers and not repeated_footers:
        return pages

    cleaned: list[PdfPage] = []
    for page in pages:
        lines = page.text.splitlines()
        removed_header_lines = 0
        removed_footer_lines = 0
        while lines and lines[0].strip().casefold() in repeated_headers:
            lines.pop(0)
            removed_header_lines += 1
        while lines and lines[-1].strip().casefold() in repeated_footers:
            lines.pop()
            removed_footer_lines += 1
        cleaned.append(
            PdfPage(
                page_number=page.page_number,
                text=normalize_text_content("\n".join(lines)),
                removed_header_lines=removed_header_lines,
                removed_footer_lines=removed_footer_lines,
            )
        )
    return cleaned


def extract_pdf_pages_local(path: Path) -> list[PdfPage]:
    reader = PdfReader(str(path))
    if reader.is_encrypted:
        try:
            decrypt_result = reader.decrypt("")
        except Exception as exc:  # pragma: no cover - library-specific failure path
            raise ValueError(f"pdf decryption failed: {exc}") from exc
        if decrypt_result == 0:
            raise ValueError("pdf is encrypted and cannot be decrypted without a password")
    pages = [
        PdfPage(page_number=page_number, text=normalize_text_content(page.extract_text() or ""))
        for page_number, page in enumerate(reader.pages, start=1)
    ]
    return [page for page in _strip_repeated_pdf_artifacts(pages) if page.text]


def _estimate_word_count(text: str) -> int:
    return len(re.findall(r"\b[\w'-]+\b", text))


def _estimate_alnum_ratio(text: str) -> float:
    visible_chars = [char for char in text if not char.isspace()]
    if not visible_chars:
        return 0.0
    alnum_chars = sum(1 for char in visible_chars if char.isalnum())
    return round(alnum_chars / len(visible_chars), 3)


def assess_pdf_local_quality(pages: list[PdfPage]) -> dict[str, Any]:
    if not pages:
        return {
            "attempted": True,
            "succeeded": True,
            "trustworthy": False,
            "page_count": 0,
            "total_text_chars": 0,
            "total_tokens": 0,
            "substantive_pages": 0,
            "substantive_page_numbers": [],
            "low_signal_pages": 0,
            "low_signal_page_numbers": [],
            "garbage_pages": 0,
            "garbage_page_numbers": [],
            "cleanup_page_numbers": [],
            "repeated_header_lines_removed": 0,
            "repeated_footer_lines_removed": 0,
            "artifact_cleanup_applied": False,
            "fallback_reasons": ["no_extractable_pages"],
            "error": None,
        }

    total_text_chars = 0
    total_tokens = 0
    substantive_page_numbers: list[int] = []
    low_signal_page_numbers: list[int] = []
    garbage_page_numbers: list[int] = []
    cleanup_page_numbers: list[int] = []
    repeated_header_lines_removed = 0
    repeated_footer_lines_removed = 0

    for page in pages:
        char_count = len(page.text)
        token_count = estimate_token_count(page.text)
        word_count = _estimate_word_count(page.text)
        alnum_ratio = _estimate_alnum_ratio(page.text)
        total_text_chars += char_count
        total_tokens += token_count
        repeated_header_lines_removed += page.removed_header_lines
        repeated_footer_lines_removed += page.removed_footer_lines
        if page.removed_header_lines or page.removed_footer_lines:
            cleanup_page_numbers.append(page.page_number)

        substantive = (
            char_count >= PDF_LOCAL_SUBSTANTIVE_CHAR_FLOOR
            and token_count >= PDF_LOCAL_SUBSTANTIVE_TOKEN_FLOOR
            and word_count >= PDF_LOCAL_SUBSTANTIVE_WORD_FLOOR
        )
        low_signal = (
            char_count < PDF_LOCAL_LOW_SIGNAL_CHAR_FLOOR
            or token_count < PDF_LOCAL_LOW_SIGNAL_TOKEN_FLOOR
            or word_count < PDF_LOCAL_LOW_SIGNAL_WORD_FLOOR
            or alnum_ratio < PDF_LOCAL_LOW_SIGNAL_ALNUM_RATIO
        )
        garbage = (
            char_count < PDF_LOCAL_GARBAGE_CHAR_FLOOR
            or token_count < PDF_LOCAL_GARBAGE_TOKEN_FLOOR
            or word_count < PDF_LOCAL_GARBAGE_WORD_FLOOR
            or alnum_ratio < PDF_LOCAL_GARBAGE_ALNUM_RATIO
        )
        if substantive:
            substantive_page_numbers.append(page.page_number)
        if low_signal:
            low_signal_page_numbers.append(page.page_number)
        if garbage:
            garbage_page_numbers.append(page.page_number)

    fallback_reasons: list[str] = []
    if total_text_chars < PDF_LOCAL_TOTAL_CHAR_FLOOR:
        fallback_reasons.append("total_text_chars_below_floor")
    if total_tokens < PDF_LOCAL_TOTAL_TOKEN_FLOOR:
        fallback_reasons.append("total_tokens_below_floor")
    if not substantive_page_numbers:
        fallback_reasons.append("no_substantive_pages")
    if len(pages) > 1 and len(low_signal_page_numbers) == len(pages):
        fallback_reasons.append("all_pages_low_signal")
    elif len(pages) >= 3 and len(low_signal_page_numbers) > len(substantive_page_numbers):
        fallback_reasons.append("low_signal_pages_dominate")
    if len(garbage_page_numbers) >= max(1, (len(pages) + 1) // 2):
        fallback_reasons.append("garbage_pages_dominate")

    return {
        "attempted": True,
        "succeeded": True,
        "trustworthy": not fallback_reasons,
        "page_count": len(pages),
        "total_text_chars": total_text_chars,
        "total_tokens": total_tokens,
        "substantive_pages": len(substantive_page_numbers),
        "substantive_page_numbers": substantive_page_numbers,
        "low_signal_pages": len(low_signal_page_numbers),
        "low_signal_page_numbers": low_signal_page_numbers,
        "garbage_pages": len(garbage_page_numbers),
        "garbage_page_numbers": garbage_page_numbers,
        "cleanup_page_numbers": cleanup_page_numbers,
        "repeated_header_lines_removed": repeated_header_lines_removed,
        "repeated_footer_lines_removed": repeated_footer_lines_removed,
        "artifact_cleanup_applied": bool(cleanup_page_numbers),
        "fallback_reasons": fallback_reasons,
        "error": None,
    }


def _local_ocr_available() -> bool:
    return bool(
        importlib.util.find_spec("rapidocr_onnxruntime")
        and importlib.util.find_spec("pypdfium2")
    )


def extract_pdf_pages_ocr_local(path: Path) -> list[PdfPage]:
    try:
        import numpy as np
        import pypdfium2 as pdfium
        from rapidocr_onnxruntime import RapidOCR
    except ImportError as exc:
        raise ValueError("local OCR dependencies are not installed") from exc

    engine = RapidOCR()
    document = pdfium.PdfDocument(str(path))
    pages: list[PdfPage] = []
    try:
        for page_number in range(len(document)):
            page = document[page_number]
            bitmap = page.render(scale=2)
            image = bitmap.to_numpy() if hasattr(bitmap, "to_numpy") else np.asarray(bitmap.to_pil())
            result, _ = engine(image)
            lines: list[str] = []
            for item in result or []:
                if isinstance(item, (list, tuple)) and len(item) >= 2 and isinstance(item[1], str):
                    lines.append(item[1])
            text = normalize_text_content("\n".join(lines))
            if text:
                pages.append(PdfPage(page_number=page_number + 1, text=text))
    finally:
        document.close()
    return pages


def build_pdf_documents(
    *,
    page_texts: list[PdfPage],
    base_metadata: dict[str, Any],
) -> list[Document]:
    documents: list[Document] = []
    for page in page_texts:
        documents.append(
            Document(
                text=page.text,
                metadata={
                    **base_metadata,
                    "page_start": page.page_number,
                    "page_end": page.page_number,
                    "section_title": detect_section_title(page.text),
                    "token_count": estimate_token_count(page.text),
                },
            )
        )
    return documents


def build_pdf_document_from_markdown(*, text: str, base_metadata: dict[str, Any]) -> list[Document]:
    cleaned = normalize_text_content(text)
    if not cleaned:
        return []
    return [
        Document(
            text=cleaned,
            metadata={
                **base_metadata,
                "page_start": 1,
                "page_end": None,
                "section_title": detect_section_title(cleaned),
                "token_count": estimate_token_count(cleaned),
            },
        )
    ]


def build_pdf_nodes(*, documents: list[Document], window_size: int) -> list[Any]:
    pipeline = IngestionPipeline(
        transformations=[
            SentenceWindowNodeParser.from_defaults(
                window_size=window_size,
                window_metadata_key=PDF_WINDOW_METADATA_KEY,
                original_text_metadata_key=PDF_ORIGINAL_TEXT_METADATA_KEY,
            )
        ]
    )
    return list(pipeline.run(documents=documents, show_progress=False))


def _empty_local_quality_report() -> dict[str, Any]:
    return {
        "attempted": False,
        "succeeded": False,
        "trustworthy": False,
        "page_count": 0,
        "total_text_chars": 0,
        "total_tokens": 0,
        "substantive_pages": 0,
        "substantive_page_numbers": [],
        "low_signal_pages": 0,
        "low_signal_page_numbers": [],
        "garbage_pages": 0,
        "garbage_page_numbers": [],
        "cleanup_page_numbers": [],
        "repeated_header_lines_removed": 0,
        "repeated_footer_lines_removed": 0,
        "artifact_cleanup_applied": False,
        "fallback_reasons": [],
        "error": None,
    }


def _empty_cloud_diagnostics(*, available: bool) -> dict[str, Any]:
    return {
        "available": available,
        "attempted": False,
        "succeeded": False,
        "parse_lane": None,
        "document_count": 0,
        "total_text_chars": 0,
        "error": None,
    }


def _empty_ocr_diagnostics(*, available: bool) -> dict[str, Any]:
    return {
        "available": available,
        "attempted": False,
        "succeeded": False,
        "trustworthy": False,
        "parse_lane": None,
        "page_count": 0,
        "total_text_chars": 0,
        "fallback_reasons": [],
        "error": None,
    }


def _format_pdf_reasons(reasons: list[str]) -> str:
    return ", ".join(reasons) if reasons else "no detailed reason recorded"


def extract_pdf_documents(
    *,
    path: Path,
    requested_lane: str,
    tier: str,
    trace_id: str,
    base_metadata: dict[str, Any],
    parse_lane_policy: str,
) -> PdfExtractionResult:
    requested_lane = requested_lane or "default"
    if requested_lane not in {"default", "local", "cloud"}:
        raise ValueError(f"unsupported requested pdf lane: {requested_lane}")
    if parse_lane_policy not in {"local_default", "cloud_default", "auto"}:
        raise ValueError(f"unsupported pdf parse lane policy: {parse_lane_policy}")

    cloud_available = bool(settings.llama_cloud_api_key)
    diagnostics = {
        "requested_lane": requested_lane,
        "parse_lane_policy": parse_lane_policy,
        "decision": None,
        "fallback_reasons": [],
        "selected_parse_lane": None,
        "local": _empty_local_quality_report(),
        "ocr": _empty_ocr_diagnostics(available=_local_ocr_available()),
        "cloud": _empty_cloud_diagnostics(available=cloud_available),
    }
    local_pages: list[PdfPage] = []
    ocr_pages: list[PdfPage] = []

    def local_failure_message(prefix: str, local_diag: dict[str, Any]) -> ValueError:
        reason_text = _format_pdf_reasons(list(local_diag.get("fallback_reasons") or ["local_parse_error"]))
        if local_diag.get("error"):
            reason_text = f"{reason_text}; {local_diag['error']}"
        return ValueError(f"{prefix}: {reason_text}")

    def attempt_local() -> tuple[list[PdfPage], dict[str, Any]]:
        nonlocal local_pages
        local_diag = diagnostics["local"]
        if local_diag["attempted"]:
            return local_pages, local_diag
        local_diag["attempted"] = True
        try:
            local_pages = extract_pdf_pages_local(path)
        except Exception as exc:
            diagnostics["local"] = {
                **_empty_local_quality_report(),
                "attempted": True,
                "succeeded": False,
                "fallback_reasons": ["local_parse_error"],
                "error": str(exc),
            }
            return [], diagnostics["local"]
        diagnostics["local"] = assess_pdf_local_quality(local_pages)
        return local_pages, diagnostics["local"]

    def attempt_ocr() -> tuple[list[PdfPage], dict[str, Any]]:
        nonlocal ocr_pages
        ocr_diag = diagnostics["ocr"]
        if ocr_diag["attempted"]:
            return ocr_pages, ocr_diag
        ocr_diag["attempted"] = True
        try:
            ocr_pages = extract_pdf_pages_ocr_local(path)
        except Exception as exc:
            diagnostics["ocr"] = {
                **_empty_ocr_diagnostics(available=_local_ocr_available()),
                "attempted": True,
                "succeeded": False,
                "fallback_reasons": ["ocr_parse_error"],
                "error": str(exc),
            }
            return [], diagnostics["ocr"]
        quality = assess_pdf_local_quality(ocr_pages)
        diagnostics["ocr"] = {
            **quality,
            "available": True,
            "attempted": True,
            "succeeded": bool(ocr_pages),
            "parse_lane": "local_ocr",
        }
        return ocr_pages, diagnostics["ocr"]

    def build_page_result(
        *,
        decision: str,
        pages: list[PdfPage],
        page_diag: dict[str, Any],
        parse_lane: str,
        fallback_reasons: list[str] | None = None,
    ) -> PdfExtractionResult:
        documents = build_pdf_documents(page_texts=pages, base_metadata=base_metadata)
        total_text_chars = sum(len(document.text) for document in documents)
        if total_text_chars == 0:
            raise local_failure_message(
                "local pdf parse produced no extractable text",
                {
                    **page_diag,
                    "fallback_reasons": list(page_diag.get("fallback_reasons") or []) + ["no_extractable_text_after_document_build"],
                },
            )
        diagnostics["decision"] = decision
        diagnostics["fallback_reasons"] = list(fallback_reasons or [])
        diagnostics["selected_parse_lane"] = parse_lane
        return PdfExtractionResult(
            documents=documents,
            parse_lane=parse_lane,
            page_count=page_diag["page_count"],
            total_text_chars=total_text_chars,
            parse_diagnostics=diagnostics,
        )

    def attempt_local_then_ocr(*, prefix: str) -> tuple[list[PdfPage], dict[str, Any], str]:
        pages, local_diag = attempt_local()
        if local_diag["trustworthy"]:
            return pages, local_diag, "local_pypdf"
        ocr_result_pages, ocr_diag = attempt_ocr()
        if ocr_diag["trustworthy"]:
            return ocr_result_pages, ocr_diag, "local_ocr"
        raise ValueError(
            f"{prefix}: "
            f"{_format_pdf_reasons(list(local_diag.get('fallback_reasons') or []))}"
            + (
                f"; ocr fallback: {_format_pdf_reasons(list(ocr_diag.get('fallback_reasons') or []))}"
                if ocr_diag.get("attempted")
                else ""
            )
            + (f"; ocr error: {ocr_diag.get('error')}" if ocr_diag.get("error") else "")
        )

    def build_cloud_result(
        *,
        decision: str,
        fallback_reasons: list[str] | None = None,
        page_count_hint: int = 0,
    ) -> PdfExtractionResult:
        cloud_diag = diagnostics["cloud"]
        cloud_diag["attempted"] = True
        text, parse_lane = extract_text_cloud(path, tier, trace_id)
        documents = build_pdf_document_from_markdown(text=text, base_metadata=base_metadata)
        total_text_chars = sum(len(document.text) for document in documents)
        if total_text_chars == 0:
            raise ValueError("cloud pdf parse returned no extractable text")
        cloud_diag.update(
            {
                "succeeded": True,
                "parse_lane": parse_lane,
                "document_count": len(documents),
                "total_text_chars": total_text_chars,
                "error": None,
            }
        )
        diagnostics["decision"] = decision
        diagnostics["fallback_reasons"] = list(fallback_reasons or [])
        diagnostics["selected_parse_lane"] = parse_lane
        return PdfExtractionResult(
            documents=documents,
            parse_lane=parse_lane,
            page_count=max(page_count_hint, 1 if documents else 0),
            total_text_chars=total_text_chars,
            parse_diagnostics=diagnostics,
        )

    if requested_lane == "cloud":
        if not cloud_available:
            diagnostics["cloud"]["error"] = "LLAMA_CLOUD_API_KEY is not set; cannot use cloud parse lane"
            raise ValueError("cloud pdf parse requested but LLAMA_CLOUD_API_KEY is not set")
        try:
            return build_cloud_result(decision="cloud_only")
        except Exception as exc:
            diagnostics["cloud"]["error"] = str(exc)
            raise

    if requested_lane == "local":
        pages, page_diag, parse_lane = attempt_local_then_ocr(prefix="local pdf parse unusable")
        return build_page_result(decision="local_only", pages=pages, page_diag=page_diag, parse_lane=parse_lane)

    if parse_lane_policy == "local_default":
        pages, page_diag, parse_lane = attempt_local_then_ocr(prefix="local-default pdf parse unusable")
        return build_page_result(
            decision="default_local_default_local_only",
            pages=pages,
            page_diag=page_diag,
            parse_lane=parse_lane,
        )

    if parse_lane_policy == "cloud_default":
        cloud_fallback_reasons: list[str] = []
        if cloud_available:
            try:
                return build_cloud_result(decision="default_cloud_default_cloud_first")
            except Exception as exc:
                diagnostics["cloud"]["error"] = str(exc)
                cloud_fallback_reasons = ["cloud_parse_failed"]
        else:
            diagnostics["cloud"]["error"] = "LLAMA_CLOUD_API_KEY is not set; cannot use cloud parse lane"
            cloud_fallback_reasons = ["cloud_unavailable"]

        pages, page_diag, parse_lane = attempt_local_then_ocr(
            prefix="cloud-default pdf parse could not stay on cloud and local fallback was unusable"
        )
        return build_page_result(
            decision="default_cloud_default_local_fallback",
            pages=pages,
            page_diag=page_diag,
            parse_lane=parse_lane,
            fallback_reasons=cloud_fallback_reasons,
        )

    pages, local_diag = attempt_local()
    if local_diag["trustworthy"]:
        return build_page_result(
            decision="default_auto_local_accepted",
            pages=pages,
            page_diag=local_diag,
            parse_lane="local_pypdf",
        )

    ocr_result_pages, ocr_diag = attempt_ocr()
    if ocr_diag["trustworthy"]:
        return build_page_result(
            decision="default_auto_local_ocr_fallback",
            pages=ocr_result_pages,
            page_diag=ocr_diag,
            parse_lane="local_ocr",
            fallback_reasons=list(local_diag.get("fallback_reasons") or []),
        )

    if cloud_available:
        try:
            return build_cloud_result(
                decision="default_auto_cloud_fallback",
                fallback_reasons=list(local_diag.get("fallback_reasons") or []),
                page_count_hint=local_diag["page_count"],
            )
        except Exception as exc:
            diagnostics["cloud"]["error"] = str(exc)
            raise ValueError(
                "auto pdf parse rejected local extraction and cloud fallback failed: "
                f"{_format_pdf_reasons(list(local_diag.get('fallback_reasons') or []))}"
                + (
                    f"; ocr fallback: {_format_pdf_reasons(list(ocr_diag.get('fallback_reasons') or []))}"
                    if ocr_diag.get("attempted")
                    else ""
                )
                + (f"; ocr error: {ocr_diag.get('error')}" if ocr_diag.get("error") else "")
                + f"; cloud error: {exc}"
            ) from exc

    diagnostics["cloud"]["error"] = "LLAMA_CLOUD_API_KEY is not set; cannot use cloud parse lane"
    raise ValueError(
        "auto pdf parse rejected local extraction and cloud fallback is unavailable: "
        f"{_format_pdf_reasons(list(local_diag.get('fallback_reasons') or []))}"
        + (
            f"; ocr fallback: {_format_pdf_reasons(list(ocr_diag.get('fallback_reasons') or []))}"
            if ocr_diag.get("attempted")
            else ""
        )
        + (f"; ocr error: {ocr_diag.get('error')}" if ocr_diag.get("error") else "")
    )


def extract_text_local(path: Path) -> tuple[str, str]:
    """Return (text, parse_lane_label) for on-box parsing."""
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md", ".csv", ".json"}:
        return normalize_text_content(path.read_text(encoding="utf-8", errors="replace")), "local_text"

    if suffix == ".pdf":
        return "\n\n".join(page.text for page in extract_pdf_pages_local(path)), "local_pypdf"

    if suffix == ".docx":
        from docx import Document as DocxDocument

        doc = DocxDocument(str(path))
        return normalize_text_content("\n".join(p.text for p in doc.paragraphs if p.text.strip())), "local_docx"

    if suffix == ".pptx":
        prs = Presentation(str(path))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
        return normalize_text_content("\n".join(lines)), "local_pptx"

    if suffix in {".xlsx", ".xlsm"}:
        structure = extract_spreadsheet_structure(path) or {"sheets": []}
        lines: list[str] = []
        for sheet in structure["sheets"]:
            lines.append(f"## Sheet: {sheet['title']}")
            for row in sheet["rows"]:
                lines.append("\t".join(row))
        return normalize_text_content("\n".join(lines)), "local_xlsx"

    if suffix in {".html", ".htm"}:
        html = path.read_text(encoding="utf-8", errors="replace")
        return normalize_text_content(BeautifulSoup(html, "html.parser").get_text(separator="\n")), "local_html"

    return normalize_text_content(path.read_text(encoding="utf-8", errors="replace")), "local_fallback"


def extract_text_cloud(path: Path, tier: str, trace_id: str) -> tuple[str, str]:
    """Parse via LlamaParse (requires LLAMA_CLOUD_API_KEY)."""
    if not settings.llama_cloud_api_key:
        raise ValueError("LLAMA_CLOUD_API_KEY is not set; cannot use cloud parse lane")

    from llama_parse import LlamaParse

    def _run() -> tuple[str, str]:
        parser = LlamaParse(
            api_key=settings.llama_cloud_api_key,
            result_type="markdown",
            verbose=False,
        )
        _ = tier  # reserved for future SDK tier flags
        documents = parser.load_data(str(path))
        text = "\n\n".join(getattr(document, "text", "") or "" for document in documents)
        return normalize_text_content(text), "llamaparse"

    return wrap_outbound_call(
        trace_id=trace_id,
        service="ghostdash-worker",
        route="llamaparse.parse",
        fn=_run,
    )
